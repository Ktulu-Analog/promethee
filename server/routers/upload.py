# ============================================================================
# Prométhée — Assistant IA avancé
# ============================================================================
# Auteur  : Pierre COUGET ktulu.analog@gmail.com
# Licence : GNU Affero General Public License v3.0 (AGPL-3.0)
#           https://www.gnu.org/licenses/agpl-3.0.html
# Année   : 2026
# ----------------------------------------------------------------------------
# Ce fichier fait partie du projet Prométhée.
# Vous pouvez le redistribuer et/ou le modifier selon les termes de la
# licence AGPL-3.0 publiée par la Free Software Foundation.
# ============================================================================

"""
routers/upload.py — Endpoint d'upload et d'extraction de fichiers

Reproduit exactement la logique de FileProcessingWorker (Qt) côté serveur :
  - Texte brut (.txt, .md, .py, .js, .json, .xml, .yaml, .csv, .html, .css)
  - PDF  → extraction texte via PyMuPDF (fitz)
  - Office → extraction via python-docx / openpyxl / python-pptx
  - Image → encodage base64 PNG (redimensionnée si nécessaire)
  - Binaires → chemin uniquement (non utilisé dans le contexte web)

Nouveauté v2 (VFS) :
  Chaque fichier uploadé est automatiquement persisté dans le système de
  fichiers virtuel de l'utilisateur sous /uploads/<nom_fichier>.
  La réponse JSON inclut le champ `vfs_path` indiquant où retrouver le
  fichier original (bytes bruts) pour les outils LLM.

Route :
    POST /upload/file
        body : multipart/form-data  { file: <UploadFile> }
        retour JSON :
            {
              type: "file"|"image",
              name: str,
              content?: str,        # contenu extrait (texte)
              base64?: str,         # image encodée base64
              mime_type?: str,
              pages?: int,          # PDF uniquement
              vfs_path: str,        # chemin dans le VFS de l'utilisateur
            }
"""

import base64
import io
import logging
import mimetypes
from pathlib import Path

from fastapi import APIRouter, UploadFile, File, HTTPException, Depends
from server.deps import require_auth
from server.routers.vfs_router import persist_upload              # ← VFS

_log = logging.getLogger(__name__)
router = APIRouter(prefix="/upload", tags=["upload"])

# Extensions texte traitées directement
_TEXT_SUFFIXES = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".json", ".xml", ".yaml", ".yml", ".csv", ".html", ".css",
    ".sh", ".bash", ".toml", ".ini", ".cfg", ".log",
}

# Limite de contenu texte (chars) — identique à FileProcessingWorker
_MAX_CONTENT_CHARS = 100_000


@router.post("/file")
async def upload_file(
    file: UploadFile = File(...),
    user: dict = Depends(require_auth),
):
    """
    Reçoit un fichier, extrait son contenu lisible et le renvoie au frontend.

    Le frontend n'a plus besoin de lire le fichier lui-même — toute l'extraction
    se fait ici, exactement comme dans FileProcessingWorker.py côté Qt.

    Le fichier est également persisté dans le VFS de l'utilisateur sous
    /uploads/<nom_fichier> pour être accessible aux outils LLM via vfs_read_file.
    """
    name = file.filename or "fichier"
    suffix = Path(name).suffix.lower()
    raw = await file.read()

    # Détection du mime type pour la persistance VFS
    mime_type, _ = mimetypes.guess_type(name)
    mime_type = mime_type or "application/octet-stream"

    # ── Texte brut ─────────────────────────────────────────────────────────
    if suffix in _TEXT_SUFFIXES:
        try:
            content = raw.decode("utf-8", errors="replace")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Lecture texte échouée : {e}")
        if len(content) > _MAX_CONTENT_CHARS:
            content = content[:_MAX_CONTENT_CHARS] + "\n\n... [tronqué]"

        vfs_path = persist_upload(user["id"], name, raw, mime_type or "text/plain")
        return {"type": "file", "name": name, "content": content,
                "vfs_path": vfs_path}

    # ── PDF ────────────────────────────────────────────────────────────────
    if suffix == ".pdf":
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise HTTPException(
                status_code=501,
                detail="PyMuPDF non installé sur le serveur (pip install pymupdf)"
            )
        try:
            doc = fitz.open(stream=raw, filetype="pdf")
            pages = len(doc)
            content_parts = []
            for i, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    content_parts.append(f"--- Page {i+1} ---\n{text}")
                if i >= 49:  # limite 50 pages max
                    content_parts.append(
                        f"[... {pages - 50} page(s) supplémentaire(s) non incluses]"
                    )
                    break
            doc.close()
            content = "\n\n".join(content_parts)
            if not content.strip():
                content = "[PDF sans texte extractible — document scanné ou image uniquement]"
            if len(content) > _MAX_CONTENT_CHARS:
                content = content[:_MAX_CONTENT_CHARS] + "\n\n... [tronqué]"

            vfs_path = persist_upload(user["id"], name, raw, "application/pdf")
            return {"type": "file", "name": name, "content": content,
                    "pages": pages, "vfs_path": vfs_path}
        except Exception as e:
            _log.exception("Extraction PDF échouée pour %s", name)
            raise HTTPException(status_code=500, detail=f"Extraction PDF échouée : {e}")

    # ── Word (.docx) ───────────────────────────────────────────────────────
    if suffix == ".docx":
        try:
            import docx
        except ImportError:
            raise HTTPException(
                status_code=501,
                detail="python-docx non installé (pip install python-docx)"
            )
        try:
            doc = docx.Document(io.BytesIO(raw))
            content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if len(content) > _MAX_CONTENT_CHARS:
                content = content[:_MAX_CONTENT_CHARS] + "\n\n... [tronqué]"

            vfs_path = persist_upload(
                user["id"], name, raw,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
            return {"type": "file", "name": name, "content": content,
                    "vfs_path": vfs_path}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Extraction Word échouée : {e}")

    # ── Excel (.xlsx) ──────────────────────────────────────────────────────
    if suffix in (".xlsx", ".xls"):
        try:
            import openpyxl
        except ImportError:
            raise HTTPException(
                status_code=501,
                detail="openpyxl non installé (pip install openpyxl)"
            )
        try:
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"=== Feuille : {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_str = "\t".join("" if v is None else str(v) for v in row)
                    if row_str.strip():
                        parts.append(row_str)
            wb.close()
            content = "\n".join(parts)
            if len(content) > _MAX_CONTENT_CHARS:
                content = content[:_MAX_CONTENT_CHARS] + "\n\n... [tronqué]"

            xls_mime = (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                if suffix == ".xlsx" else "application/vnd.ms-excel"
            )
            vfs_path = persist_upload(user["id"], name, raw, xls_mime)
            return {"type": "file", "name": name, "content": content,
                    "vfs_path": vfs_path}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Extraction Excel échouée : {e}")

    # ── PowerPoint (.pptx) ─────────────────────────────────────────────────
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise HTTPException(
                status_code=501,
                detail="python-pptx non installé (pip install python-pptx)"
            )
        try:
            prs = Presentation(io.BytesIO(raw))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text_frame.text
                    for shape in slide.shapes
                    if hasattr(shape, "text_frame") and shape.text_frame.text.strip()
                ]
                if texts:
                    parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
            content = "\n\n".join(parts)
            if len(content) > _MAX_CONTENT_CHARS:
                content = content[:_MAX_CONTENT_CHARS] + "\n\n... [tronqué]"

            vfs_path = persist_upload(
                user["id"], name, raw,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            return {"type": "file", "name": name, "content": content,
                    "vfs_path": vfs_path}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Extraction PowerPoint échouée : {e}")

    # ── Image ──────────────────────────────────────────────────────────────
    if suffix in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
        img_mime = {
            ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".png": "image/png",  ".gif": "image/gif",
            ".webp": "image/webp", ".bmp": "image/bmp",
        }.get(suffix, "image/png")

        try:
            from PIL import Image as PILImage
            img = PILImage.open(io.BytesIO(raw))
            # Redimensionner si trop grande (identique à FileProcessingWorker)
            max_size = 1024
            if img.width > max_size or img.height > max_size:
                img.thumbnail((max_size, max_size), PILImage.LANCZOS)
            # Convertir en PNG base64
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            width, height = img.width, img.height
        except ImportError:
            # Fallback sans Pillow : encoder les bytes bruts tels quels
            _log.warning("Pillow non installé — fallback base64 brut pour %s", name)
            b64 = base64.b64encode(raw).decode("ascii")
            width = height = None

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Traitement image échoué : {e}")

        # Persister l'image originale (bytes bruts) dans le VFS
        vfs_path = persist_upload(user["id"], name, raw, img_mime)

        result = {
            "type": "image", "name": name,
            "base64": b64, "mime_type": "image/png",
            "vfs_path": vfs_path,
        }
        if width is not None:
            result["width"] = width
            result["height"] = height
        return result

    # ── Fichier binaire non reconnu ────────────────────────────────────────
    vfs_path = persist_upload(user["id"], name, raw, "application/octet-stream")
    return {
        "type": "file",
        "name": name,
        "content": f"[Fichier binaire : {suffix} — {len(raw)} octets]",
        "vfs_path": vfs_path,
    }
