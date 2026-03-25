# ============================================================================
# Prométhée — Assistant IA desktop
# ============================================================================
# Auteur  : Pierre COUGET
# Licence : GNU Affero General Public License v3.0 (AGPL-3.0)
#           https://www.gnu.org/licenses/agpl-3.0.html
# Année   : 2026
# ----------------------------------------------------------------------------
# Ce fichier fait partie du projet Prométhée.
# Vous pouvez le redistribuer et/ou le modifier selon les termes de la
# licence AGPL-3.0 publiée par la Free Software Foundation.
# ============================================================================

"""
tools/export_tools.py — Génération et export de fichiers bureautiques
======================================================================

Outils exposés (9) :

  Markdown (1) :
    - export_md         : écrit un fichier Markdown à partir de contenu texte

  Word / docx (1) :
    - export_docx       : génère un document Word structuré (titres, paragraphes,
                          tableaux, listes) depuis une description JSON

  Tableur / xlsx (2) :
    - export_xlsx_json  : génère un classeur Excel depuis une structure JSON
                          (feuilles, en-têtes, lignes, graphiques)
    - export_xlsx_csv   : génère un classeur Excel depuis du CSV brut (une feuille)

  Présentation / pptx (2) :
    - export_pptx_json  : génère une présentation PowerPoint depuis une structure
                          JSON (titre, slides, puces, notes)
    - export_pptx_outline: génère une présentation depuis un outline texte
                          (ligne "# Titre", "- Puce", "> Note")

  PDF (1) :
    - export_pdf        : génère un PDF structuré (titres, paragraphes, tableaux)
                          depuis une description JSON, via reportlab

  LibreOffice natif (2) :
    - export_libreoffice: convertit un fichier existant vers odt/ods/odp
                          en invoquant LibreOffice headless
    - export_libreoffice_native : génère directement un odt/ods/odp depuis
                          une description JSON (via python-docx/openpyxl/python-pptx
                          + conversion LibreOffice)

Conventions communes
────────────────────
  - output_path : chemin absolu ou relatif au home utilisateur.
                  Si omis ou vide, un fichier est créé dans ~/Exports/Prométhée/.
  - Retour      : dict JSON {"path": "/chemin/absolu", "size_bytes": N,
                             "pages"/"sheets"/"slides": N, "status": "ok"}
  - En cas d'erreur : {"error": "message explicatif", "status": "error"}

Structure JSON commune pour export_docx, export_pdf, export_pptx_json
──────────────────────────────────────────────────────────────────────
  {
    "title": "Titre du document",
    "sections": [
      {
        "heading": "Titre de section",   // niveau 1-3 (optionnel)
        "level": 1,                      // 1, 2 ou 3 (défaut 1)

        // Contenu textuel — trois variantes (choisir la plus adaptée) :
        "paragraphs": ["Para 1 développé.", "Para 2 développé."],
                                         // PRÉFÉRÉ pour plusieurs paragraphes
        "content": "Texte avec\n\ndoubles sauts\n\nou simple paragraphe",
                                         // alternatif à paragraphs
        "intro": "Texte avant bullets/tableau",  // cumulable avec bullets/table

        // Éléments structurés (cumulables avec intro) :
        "table": {
          "headers": ["Col A", "Col B"],
          "rows": [["val1", "val2"], ["val3", "val4"]]
        },
        "bullets": ["item 1", "item 2"],

        "page_break": false              // saut de page optionnel après la section
      }
    ]
  }

  NOTE : un document professionnel complet comporte typiquement 15 à 40 sections.
  Chaque section doit contenir un contenu rédigé, dense et proportionnel au sujet.

Prérequis :
    pip install python-docx openpyxl python-pptx reportlab
    LibreOffice installé système (apt install libreoffice)
"""

import io
import json
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

from core.tools_engine import tool, set_current_family

set_current_family("export_tools", "Export de fichiers", "📄")

# ── Répertoire de sortie par défaut ──────────────────────────────────────────

_DEFAULT_EXPORT_DIR = Path.home() / "Exports" / "Prométhée"


def _resolve_output(output_path: str, default_name: str) -> Path:
    """Résout le chemin de sortie et crée les répertoires manquants."""
    if output_path:
        p = Path(output_path).expanduser()
        if not p.is_absolute():
            p = Path.home() / p
    else:
        _DEFAULT_EXPORT_DIR.mkdir(parents=True, exist_ok=True)
        p = _DEFAULT_EXPORT_DIR / default_name
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _ok(path: Path, extra: dict | None = None) -> str:
    r = {"status": "ok", "path": str(path), "size_bytes": path.stat().st_size}
    if extra:
        r.update(extra)
    return json.dumps(r, ensure_ascii=False)


def _err(msg: str) -> str:
    return json.dumps({"status": "error", "error": msg}, ensure_ascii=False)


# ═════════════════════════════════════════════════════════════════════════════
# MARKDOWN
# ═════════════════════════════════════════════════════════════════════════════

@tool(
    name="export_md",
    description=(
        "Écrit un fichier Markdown (.md) à partir de contenu texte brut. "
        "À utiliser quand l'utilisateur demande de produire un rapport, un README, "
        "un article ou tout document texte en format Markdown. "
        "Le contenu doit déjà être en Markdown valide (titres #, listes -, tableaux |)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {
                "type": "string",
                "description": "Contenu Markdown complet à écrire dans le fichier."
            },
            "output_path": {
                "type": "string",
                "description": (
                    "Chemin de destination, ex: ~/Documents/rapport.md. "
                    "Si omis, crée le fichier dans ~/Exports/Prométhée/."
                )
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis (ex: rapport.md)."
            }
        },
        "required": ["content"]
    }
)
def export_md(content: str, output_path: str = "", filename: str = "") -> str:
    try:
        name = filename or "export.md"
        if not name.endswith(".md"):
            name += ".md"
        p = _resolve_output(output_path, name)
        p.write_text(content, encoding="utf-8")
        lines = content.count("\n") + 1
        return _ok(p, {"lines": lines})
    except Exception as e:
        return _err(f"export_md : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# WORD / DOCX
# ═════════════════════════════════════════════════════════════════════════════

def _parse_inline_markup(text: str) -> list[dict]:
    """
    Découpe un texte contenant du Markdown inline et des formules LaTeX
    en une liste de segments typés. Types produits :

      {"type": "text",          "content": "..."}   texte ordinaire (peut contenir \\n)
      {"type": "bold",          "content": "..."}   **gras**
      {"type": "italic",        "content": "..."}   *italique*
      {"type": "underline",     "content": "..."}   __souligné__
      {"type": "code",          "content": "..."}   `code inline`
      {"type": "code_block",    "lang": "...",
                                "content": "..."}   ```lang\\n...\\n```
      {"type": "latex_inline",  "expr":    "..."}   $formule$
      {"type": "latex_display", "expr":    "..."}   $$formule$$

    Ordre de priorité (premier match gagne) :
      ``` > $$ > $ > ** > * > __ > `
    Les délimiteurs non fermés sont laissés intacts (texte brut).
    """
    PAT = re.compile(
        r'(```(\w*)\n(.*?)```'    # bloc code     ```lang\n...\n```   (priorité max)
        r'|\\\[.*?\\\]'           # LaTeX display \[...\]
        r'|\$\$.*?\$\$'           # LaTeX display $$...$$
        r'|\\\(.*?\\\)'           # LaTeX inline  \(...\)
        r'|\$[^\$\n]+?\$'         # LaTeX inline  $...$
        r'|\*\*(.+?)\*\*'         # gras          **...**
        r'|\*([^\*\n]+?)\*'       # italique      *...*
        r'|__([^_\n]+?)__'        # souligné      __...__
        r'|`([^`\n]+?)`'          # code inline   `...`
        r')',
        re.DOTALL,
    )
    segments: list[dict] = []
    last = 0
    for m in PAT.finditer(text):
        if m.start() > last:
            segments.append({"type": "text", "content": text[last:m.start()]})
        tok = m.group(0)
        if tok.startswith("```"):
            segments.append({
                "type":    "code_block",
                "lang":    (m.group(2) or "").strip(),
                "content": m.group(3) or "",
            })
        elif tok.startswith(r"\["):
            segments.append({"type": "latex_display", "expr": tok[2:-2].strip()})
        elif tok.startswith("$$"):
            segments.append({"type": "latex_display", "expr": tok[2:-2].strip()})
        elif tok.startswith(r"\("):
            segments.append({"type": "latex_inline",  "expr": tok[2:-2].strip()})
        elif tok.startswith("$"):
            segments.append({"type": "latex_inline",  "expr": tok[1:-1].strip()})
        elif tok.startswith("**"):
            segments.append({"type": "bold",      "content": tok[2:-2]})
        elif tok.startswith("*"):
            segments.append({"type": "italic",    "content": tok[1:-1]})
        elif tok.startswith("__"):
            segments.append({"type": "underline", "content": tok[2:-2]})
        elif tok.startswith("`"):
            segments.append({"type": "code",      "content": tok[1:-1]})
        last = m.end()
    if last < len(text):
        segments.append({"type": "text", "content": text[last:]})
    return segments


# Alias conservé pour compatibilité interne
_latex_parse_runs = _parse_inline_markup


def _has_inline_markup(text: str) -> bool:
    """Retourne True si le texte contient LaTeX ou du Markdown inline."""
    return bool(re.search(r'[\$\*_`]', text))


# Alias pour compatibilité avec le code existant
_has_latex = _has_inline_markup


def _docx_add_paragraph_latex(doc, text: str, style_name: str = "Normal",
                               font_size_pt: float = 10.0, bold: bool = False) -> list:
    """
    Ajoute un ou plusieurs paragraphes au document Word en gérant les formules LaTeX.

    - Texte pur                → un seul paragraphe, comportement identique à avant.
    - Formule inline $...$     → image PNG insérée dans le run, hauteur calée sur
                                 la taille de police (alignement vertical correct).
    - Formule display $$...$$ → paragraphe centré dédié intercalé entre le texte
                                 précédent et suivant ; image plus grande.

    Retourne la liste des paragraphes python-docx créés.
    En cas d'erreur de compilation LaTeX, insère le texte brut entre crochets.
    """
    import io as _io
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    segs = _parse_inline_markup(text)
    has_display = any(s["type"] == "latex_display" for s in segs)
    created: list = []

    def _png_inline(expr: str) -> bytes:
        return _latex_to_png_bytes(expr, dpi=200)

    def _png_display(expr: str) -> bytes:
        return _latex_to_png_bytes(expr, dpi=220)

    def _flush_inline(p, inline_segs: list, fs: float, b: bool):
        """Ajoute une séquence de segments inline dans le paragraphe p."""
        for s in inline_segs:
            t = s["type"]
            if t == "text":
                # Préserver les \n comme sauts de ligne dans le paragraphe Word
                content = s["content"]
                parts = content.split("\n")
                for i, part in enumerate(parts):
                    if part:
                        run = p.add_run(part)
                        run.font.size = Pt(fs)
                        run.bold = b
                    if i < len(parts) - 1:
                        p.add_run().add_break()
            elif t == "bold":
                run = p.add_run(s["content"])
                run.font.size = Pt(fs)
                run.bold = True
            elif t == "italic":
                run = p.add_run(s["content"])
                run.font.size = Pt(fs)
                run.italic = True
                run.bold = b
            elif t == "underline":
                run = p.add_run(s["content"])
                run.font.size = Pt(fs)
                run.underline = True
                run.bold = b
            elif t == "code":
                run = p.add_run(s["content"])
                run.font.size = Pt(fs * 0.88)
                run.font.name = "Courier New"
                run.bold = False
            elif t == "latex_inline":
                try:
                    png = _png_inline(s["expr"])
                    _, h_cm = _latex_png_dims(png, 200)
                    h_pt = max(fs * 0.9, min(h_cm * 28.35, fs * 1.6))
                    run = p.add_run()
                    run.add_picture(_io.BytesIO(png), height=Pt(h_pt))
                except Exception:
                    run = p.add_run(f"[{s['expr']}]")
                    run.font.size = Pt(fs)
                    run.bold = b

    def _add_code_block(seg: dict):
        """Insère un bloc de code comme paragraphe(s) Courier New avec fond gris."""
        from docx.oxml.ns import qn as _qn
        from docx.oxml import OxmlElement as _OxmlEl
        from docx.shared import RGBColor as _RGB

        lines = seg["content"].rstrip("\n").split("\n")
        lang  = seg.get("lang", "")

        # Paragraphe d'en-tête avec le nom du langage (fond bleu, texte blanc)
        if lang:
            p_hdr = doc.add_paragraph()
            run_hdr = p_hdr.add_run(f"  {lang}")
            run_hdr.font.size = Pt(7.5)
            run_hdr.font.name = "Courier New"
            run_hdr.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF)
            run_hdr.bold = True
            # Fond bleu via shading XML
            pPr = p_hdr._p.get_or_add_pPr()
            shd = _OxmlEl("w:shd")
            shd.set(_qn("w:val"), "clear")
            shd.set(_qn("w:color"), "auto")
            shd.set(_qn("w:fill"), "2255A4")
            pPr.append(shd)
            p_hdr.paragraph_format.space_before = Pt(6)
            p_hdr.paragraph_format.space_after  = Pt(0)
            created.append(p_hdr)

        # Lignes de code (fond gris clair)
        for i, line in enumerate(lines):
            p_code = doc.add_paragraph()
            run_code = p_code.add_run(line if line else " ")
            run_code.font.size = Pt(8.5)
            run_code.font.name = "Courier New"
            run_code.font.color.rgb = _RGB(0x1a, 0x1a, 0x2e)
            pPr = p_code._p.get_or_add_pPr()
            shd = _OxmlEl("w:shd")
            shd.set(_qn("w:val"), "clear")
            shd.set(_qn("w:color"), "auto")
            shd.set(_qn("w:fill"), "F8F9FB")
            pPr.append(shd)
            p_code.paragraph_format.space_before = Pt(0)
            p_code.paragraph_format.space_after  = Pt(0)
            p_code.paragraph_format.left_indent  = Pt(10)
            created.append(p_code)

        # Espace après le bloc
        p_end = doc.add_paragraph()
        p_end.paragraph_format.space_before = Pt(0)
        p_end.paragraph_format.space_after  = Pt(4)
        created.append(p_end)

    # Types qui nécessitent un paragraphe dédié (éclater la séquence)
    BLOCK_TYPES = {"latex_display", "code_block"}

    has_blocks = any(s["type"] in BLOCK_TYPES for s in segs)

    if not has_blocks:
        # ── Cas simple : tout en ligne dans un seul paragraphe ────────────
        p = doc.add_paragraph(style=style_name)
        _flush_inline(p, segs, font_size_pt, bold)
        created.append(p)
    else:
        # ── Cas complexe : éclater aux blocs display et code ──────────────
        pending: list[dict] = []
        for seg in segs:
            if seg["type"] in BLOCK_TYPES:
                # Flush les segments inline/texte accumulés
                if pending:
                    p = doc.add_paragraph(style=style_name)
                    _flush_inline(p, pending, font_size_pt, bold)
                    created.append(p)
                    pending = []
                # Insérer le bloc dédié
                if seg["type"] == "latex_display":
                    p_disp = doc.add_paragraph()
                    p_disp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    try:
                        png = _png_display(seg["expr"])
                        _, h_cm = _latex_png_dims(png, 220)
                        h_pt = min(h_cm * 28.35, font_size_pt * 3.5)
                        run = p_disp.add_run()
                        run.add_picture(_io.BytesIO(png), height=Pt(h_pt))
                    except Exception:
                        run = p_disp.add_run(f"[{seg['expr']}]")
                        run.font.size = Pt(font_size_pt)
                    created.append(p_disp)
                elif seg["type"] == "code_block":
                    _add_code_block(seg)
            else:
                pending.append(seg)
        # Flush final
        if pending:
            p = doc.add_paragraph(style=style_name)
            _flush_inline(p, pending, font_size_pt, bold)
            created.append(p)

    return created


def _latex_to_png_bytes(expr: str, dpi: int = 200) -> bytes:
    """
    Rend une expression LaTeX en PNG bytes, puis rogne les bords transparents.

    Pipeline 1 (préféré)  : latex + dvipng   — fond transparent natif
    Pipeline 2 (repli)    : pdflatex + pdftoppm — disponible même sans dvipng ;
                            fond blanc converti en transparent par Pillow.

    Prérequis minimaux : pdflatex + pdftoppm (poppler-utils) + Pillow.
    dvipng optionnel (meilleur rendu transparent si présent).
    Lève RuntimeError si aucun pipeline ne fonctionne.
    """
    import shutil as _sh
    import subprocess as _sp
    import tempfile as _tmp
    from pathlib import Path as _P

    body = rf"\begin{{equation*}}{expr}\end{{equation*}}"
    tex_src = (
        r"\documentclass[12pt]{article}" + "\n"
        r"\usepackage{amsmath,amssymb,amsfonts,bm}" + "\n"
        r"\usepackage{xcolor}" + "\n"
        r"\usepackage[active,tightpage]{preview}" + "\n"
        r"\PreviewEnvironment{equation*}" + "\n"
        r"\setlength\PreviewBorder{2pt}" + "\n"
        r"\pagestyle{empty}" + "\n"
        r"\begin{document}" + "\n"
        r"\color[rgb]{0.102,0.102,0.122}" + "\n"
        + body + "\n"
        + r"\end{document}" + "\n"
    )

    raw_bytes: bytes | None = None
    last_error = "aucun pipeline disponible"

    with _tmp.TemporaryDirectory() as tmp:
        d = _P(tmp)
        (d / "f.tex").write_text(tex_src, encoding="utf-8")

        # ── Pipeline 1 : pdflatex + pdftoppm ─────────────────────────────
        # Prioritaire : embarque les fontes Type1 directement, pas de dépendance
        # aux fontes PK bitmap (évite le carré □ pour \Box, \square, etc.)
        if _sh.which("pdflatex") and _sh.which("pdftoppm"):
            r1 = _sp.run(
                ["pdflatex", "-interaction=nonstopmode", "-halt-on-error", "f.tex"],
                cwd=d, capture_output=True, timeout=30,
            )
            pdf_file = d / "f.pdf"
            if pdf_file.exists():
                r2 = _sp.run(
                    ["pdftoppm", "-r", str(dpi), "-png", "f.pdf", "page"],
                    cwd=d, capture_output=True, timeout=20,
                )
                # pdftoppm génère page-1.png (ou page-01.png selon la version)
                pages = sorted(d.glob("page*.png"))
                if pages:
                    raw_bytes = pages[0].read_bytes()
                    # pdftoppm produit un fond blanc — le convertir en transparent
                    try:
                        import io as _io
                        import numpy as _np
                        from PIL import Image as _Img
                        img = _Img.open(_io.BytesIO(raw_bytes)).convert("RGBA")
                        arr = _np.array(img)
                        # Pixels quasi-blancs (R,G,B > 250) → transparent
                        white = (arr[:, :, 0] > 250) & (arr[:, :, 1] > 250) & (arr[:, :, 2] > 250)
                        arr[white, 3] = 0
                        img2 = _Img.fromarray(arr, "RGBA")
                        buf = _io.BytesIO()
                        img2.save(buf, format="PNG")
                        raw_bytes = buf.getvalue()
                    except Exception:
                        pass  # Garder le PNG brut si Pillow/numpy indisponible
                else:
                    last_error = f"pdftoppm : aucun PNG généré\n{r2.stderr.decode(errors='replace')[-200:]}"
            else:
                log = (d / "f.log").read_text(errors="replace")[-500:] \
                      if (d / "f.log").exists() else ""
                last_error = f"pdflatex : compilation échouée\n{log}"

        # ── Pipeline 2 : latex + dvipng ───────────────────────────────────
        # Repli si pdflatex/pdftoppm indisponibles — fond transparent natif
        if raw_bytes is None and _sh.which("dvipng"):
            r1 = _sp.run(
                ["latex", "-interaction=nonstopmode", "-halt-on-error", "f.tex"],
                cwd=d, capture_output=True, timeout=20,
            )
            dvi = d / "f.dvi"
            if dvi.exists():
                r2 = _sp.run(
                    ["dvipng", "-D", str(dpi), "-T", "tight", "--png",
                     "-bg", "Transparent", "-o", "f.png", "f.dvi"],
                    cwd=d, capture_output=True, timeout=20,
                )
                png_file = d / "f.png"
                if png_file.exists():
                    raw_bytes = png_file.read_bytes()
                else:
                    last_error = f"dvipng : {r2.stderr.decode(errors='replace')[-300:]}"
            else:
                log = (d / "f.log").read_text(errors="replace")[-500:] \
                      if (d / "f.log").exists() else ""
                last_error = f"latex (dvi) : compilation échouée\n{log}"

    if raw_bytes is None:
        raise RuntimeError(f"_latex_to_png_bytes : {last_error}")

    # ── Crop des bords transparents ───────────────────────────────────────────
    try:
        import io as _io
        import numpy as _np
        from PIL import Image as _Img

        img  = _Img.open(_io.BytesIO(raw_bytes)).convert("RGBA")
        arr  = _np.array(img)
        alpha = arr[:, :, 3]
        rows  = _np.any(alpha > 0, axis=1)
        cols  = _np.any(alpha > 0, axis=0)
        if rows.any() and cols.any():
            PAD  = 4   # pixels de marge autour de la formule
            rmin = max(0,              _np.where(rows)[0][0]  - PAD)
            rmax = min(arr.shape[0]-1, _np.where(rows)[0][-1] + PAD)
            cmin = max(0,              _np.where(cols)[0][0]  - PAD)
            cmax = min(arr.shape[1]-1, _np.where(cols)[0][-1] + PAD)
            cropped = img.crop((cmin, rmin, cmax + 1, rmax + 1))
            buf = _io.BytesIO()
            cropped.save(buf, format="PNG")
            return buf.getvalue()
    except Exception:
        pass   # Si PIL/numpy absent ou erreur, on retourne le PNG brut

    return raw_bytes


def _latex_png_dims(png_bytes: bytes, dpi: int) -> tuple[float, float]:
    """
    Retourne les dimensions (largeur_cm, hauteur_cm) d'un PNG LaTeX
    à la résolution de rendu donnée. Utilisé pour calculer les tailles
    d'insertion dans HTML (CSS pt) et DOCX (python-docx Pt).
    """
    try:
        import io as _io
        from PIL import Image as _Img
        img = _Img.open(_io.BytesIO(png_bytes))
        w_px, h_px = img.size
        return w_px / dpi * 2.54, h_px / dpi * 2.54   # cm
    except Exception:
        return 5.0, 0.5  # fallback raisonnable


def _build_docx(doc_json: dict):
    """Construit un docx à partir de la structure JSON commune.

    Améliorations v3 :
    - 'paragraphs' : liste de chaînes → plusieurs paragraphes dans une section
    - 'content' accepte les sauts de ligne (\n) comme séparateurs de paragraphes
    - 'intro' : paragraphe introductif avant les bullets ou le tableau
    - 'bullets' et 'table' peuvent coexister avec 'content'/'intro'
    - Support LaTeX : formules $...$ (inline) et $$...$$ (display centré)
      rendues via latex + dvipng et insérées comme images PNG
    - Compatibilité totale avec les documents générés avant cette version
    """
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Titre principal
    if doc_json.get("title"):
        doc.add_heading(doc_json["title"], level=0)

    def _add_text(text: str, style_name: str = "Normal", font_size_pt: float = 10.0):
        """Ajoute un paragraphe en gérant automatiquement le LaTeX si présent."""
        text = str(text)
        if _has_latex(text):
            _docx_add_paragraph_latex(doc, text, style_name=style_name,
                                      font_size_pt=font_size_pt)
        else:
            doc.add_paragraph(text, style=style_name)

    page_count = 1
    for section in doc_json.get("sections", []):
        heading = section.get("heading")
        level   = max(1, min(3, int(section.get("level", 1))))

        if heading:
            doc.add_heading(heading, level=level)

        # ── Paragraphes ────────────────────────────────────────────────
        if section.get("paragraphs"):
            for para in section["paragraphs"]:
                if para and str(para).strip():
                    _add_text(str(para))

        elif section.get("content"):
            raw = section["content"]
            if "\n\n" in raw:
                parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
            else:
                parts = [p.strip() for p in raw.split("\n") if p.strip()]
            for part in parts:
                _add_text(part)

        # ── Intro ──────────────────────────────────────────────────────
        if section.get("intro"):
            _add_text(str(section["intro"]))

        # ── Liste à puces ───────────────────────────────────────────────
        if section.get("bullets"):
            for item in section["bullets"]:
                if item and str(item).strip():
                    _add_text(str(item), style_name="List Bullet")

        # ── Tableau ─────────────────────────────────────────────────────
        if section.get("table"):
            tbl_data = section["table"]
            headers = tbl_data.get("headers", [])
            rows    = tbl_data.get("rows", [])
            if headers:
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = "Light Grid Accent 1"
                hdr_cells = table.rows[0].cells
                for i, h in enumerate(headers):
                    hdr_cells[i].text = str(h)
                for ri, row in enumerate(rows):
                    row_cells = table.rows[ri + 1].cells
                    for ci, val in enumerate(row[:len(headers)]):
                        cell_text = str(val)
                        if _has_latex(cell_text):
                            # Vider le texte par défaut et remplir avec LaTeX
                            row_cells[ci].text = ""
                            _docx_add_paragraph_latex(
                                row_cells[ci], cell_text, font_size_pt=9.0
                            )
                        else:
                            row_cells[ci].text = cell_text
                doc.add_paragraph()  # espace après tableau

        # Page break optionnel
        if section.get("page_break"):
            doc.add_page_break()
            page_count += 1

    return doc


@tool(
    name="export_docx",
    description=(
        "Génère un document Word (.docx) structuré depuis une description JSON. "
        "Supporte les titres hiérarchiques (niveaux 1 à 3), les listes à puces, les tableaux "
        "et plusieurs paragraphes par section. "
        "À utiliser pour tout document formel : rapport, compte-rendu, contrat, note, guide. "
        "\n\n"
        "RÈGLE DE CONTENU CRITIQUE : le document doit être COMPLET et DÉVELOPPÉ. "
        "Ne pas se limiter à une structure squelette. "
        "Chaque section doit contenir un vrai contenu rédigé, proportionnel au sujet. "
        "Un rapport sérieux comporte typiquement 15 à 40 sections. "
        "Utiliser 'paragraphs' (liste) pour plusieurs paragraphes dans une section, "
        "ou 'content' avec des doubles sauts de ligne (\\n\\n) comme séparateurs. "
        "'bullets' et 'table' peuvent être combinés avec 'intro' dans la même section. "
        "\n\n"
        "Structure JSON d'une section complète : "
        '{"heading": "Titre", "level": 1, '
        '"intro": "Texte introductif (optionnel, avant les bullets/tableau)", '
        '"paragraphs": ["Paragraphe 1 développé.", "Paragraphe 2 développé."], '
        '"bullets": ["Point clé 1", "Point clé 2"], '
        '"table": {"headers": ["Col A", "Col B"], "rows": [["v1","v2"]]}, '
        '"page_break": false}'
    ),
    parameters={
        "type": "object",
        "properties": {
            "document": {
                "type": "object",
                "description": (
                    "Structure du document Word à générer. "
                    "Champ 'title' (str, titre principal). "
                    "Champ 'sections' : liste d'objets section. "
                    "Chaque section accepte : "
                    "heading (str, titre de section), "
                    "level (int 1-3, niveau de titre, défaut 1), "
                    "paragraphs (liste de str — PRÉFÉRER à 'content' pour plusieurs paragraphes), "
                    "content (str — paragraphe unique ou texte avec \\n\\n pour multi-paragraphes), "
                    "intro (str — paragraphe placé avant bullets ou tableau), "
                    "bullets (liste de str — puces, cumulable avec intro), "
                    "table (objet {headers, rows} — tableau, cumulable avec intro), "
                    "page_break (bool — saut de page après la section, défaut false). "
                    "IMPORTANT : rédiger un contenu dense et complet ; "
                    "un document de qualité professionnelle comporte 15 à 40 sections."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/rapport.docx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["document"]
    }
)
def export_docx(document: dict, output_path: str = "", filename: str = "") -> str:
    try:
        name = filename or (document.get("title", "export") + ".docx")
        if not name.endswith(".docx"):
            name += ".docx"
        p = _resolve_output(output_path, name)
        doc = _build_docx(document)
        doc.save(str(p))
        sections = len(document.get("sections", []))
        return _ok(p, {"sections": sections})
    except Exception as e:
        return _err(f"export_docx : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# EXCEL / XLSX
# ═════════════════════════════════════════════════════════════════════════════

# ── Types de graphiques supportés ────────────────────────────────────────────
#
#  bar        Barres verticales groupées       (comparaison entre catégories)
#  bar_stacked Barres verticales empilées      (composition + total)
#  bar_percent Barres empilées 100 %           (proportions)
#  bar_h       Barres horizontales groupées    (libellés longs)
#  line        Courbes                         (évolution temporelle)
#  line_smooth Courbes lissées (spline)        (tendances)
#  area        Aires                           (volumes cumulés)
#  area_stacked Aires empilées                 (parts cumulées)
#  pie         Camembert                       (parts d'un tout, ≤ 6 séries)
#  doughnut    Anneau                          (idem camembert, style moderne)
#  scatter     Nuage de points                 (corrélations entre 2 variables)
#  bubble      Bulles                          (3 variables : x, y, taille)
#  radar       Radar / toile d'araignée        (profils multi-critères)
#
# ── Structure d'un graphique dans sheets[].charts[] ─────────────────────────
#
#  {
#    "type": "bar",                  # obligatoire — voir liste ci-dessus
#    "title": "Ventes par mois",     # titre affiché au-dessus du graphique
#    "data_sheet": "Ventes",         # feuille source (défaut : feuille courante)
#    "categories_col": 1,            # colonne des étiquettes X (indice 1-based)
#    "series": [                     # liste des séries à tracer
#      {
#        "title": "CA (€)",          # légende de la série
#        "col": 2,                   # colonne des valeurs (indice 1-based)
#        "color": "2255A4"           # couleur hex optionnelle (sans #)
#      }
#    ],
#    "data_rows": [2, 13],           # [première_ligne, dernière_ligne] données
#                                    # si omis : toutes les lignes après l'en-tête
#    "anchor": "E2",                 # cellule d'ancrage du coin supérieur gauche
#                                    # (défaut : colonne après les données, ligne 1)
#    "width_cm": 15,                 # largeur en cm (défaut : 15)
#    "height_cm": 10,                # hauteur en cm (défaut : 10)
#    "style": 10,                    # style Excel 1-48 (défaut : 10)
#    "show_legend": true,            # afficher la légende (défaut : true)
#    "show_data_labels": false       # afficher les valeurs sur les barres (défaut : false)
#  }

def _build_charts(wb, ws, charts_def: list, sheet_data_rows: int, sheet_headers_count: int) -> list[str]:
    """
    Construit les graphiques openpyxl pour une feuille.
    Retourne la liste des avertissements éventuels.
    """
    from openpyxl.chart import (
        BarChart, LineChart, AreaChart, PieChart, DoughnutChart,
        ScatterChart, BubbleChart, RadarChart,
        Reference, Series,
    )
    from openpyxl.chart.series import SeriesLabel
    from openpyxl.chart.label import DataLabel
    from openpyxl.utils import get_column_letter
    import openpyxl.utils.units as units_util

    warnings = []

    _CHART_FACTORIES = {
        "bar":          lambda: _make_bar(grouping="clustered",   direction="col"),
        "bar_stacked":  lambda: _make_bar(grouping="stacked",     direction="col"),
        "bar_percent":  lambda: _make_bar(grouping="percentStacked", direction="col"),
        "bar_h":        lambda: _make_bar(grouping="clustered",   direction="bar"),
        "line":         lambda: LineChart(),
        "line_smooth":  lambda: _make_line_smooth(),
        "area":         lambda: AreaChart(),
        "area_stacked": lambda: _make_area_stacked(),
        "pie":          lambda: PieChart(),
        "doughnut":     lambda: DoughnutChart(),
        "scatter":      lambda: ScatterChart(),
        "bubble":       lambda: BubbleChart(),
        "radar":        lambda: RadarChart(),
    }

    def _make_bar(grouping: str, direction: str):
        c = BarChart()
        c.type = direction
        c.grouping = grouping
        return c

    def _make_line_smooth():
        c = LineChart()
        c.smooth = True
        return c

    def _make_area_stacked():
        c = AreaChart()
        c.grouping = "stacked"
        return c

    for idx, chart_def in enumerate(charts_def):
        chart_type = str(chart_def.get("type", "bar")).lower()
        if chart_type not in _CHART_FACTORIES:
            warnings.append(f"Graphique {idx+1} : type inconnu '{chart_type}', ignoré.")
            continue

        try:
            chart = _CHART_FACTORIES[chart_type]()

            # Titre
            if chart_def.get("title"):
                chart.title = chart_def["title"]

            # Style
            style = chart_def.get("style", 10)
            if isinstance(style, int) and 1 <= style <= 48:
                chart.style = style

            # Légende
            if not chart_def.get("show_legend", True):
                chart.legend = None

            # Dimensions
            w_cm = chart_def.get("width_cm", 15)
            h_cm = chart_def.get("height_cm", 10)
            chart.width  = w_cm
            chart.height = h_cm

            # Feuille source des données
            src_sheet_name = chart_def.get("data_sheet")
            src_ws = wb[src_sheet_name] if src_sheet_name and src_sheet_name in wb.sheetnames else ws

            # Plage de lignes des données
            data_rows = chart_def.get("data_rows")
            if data_rows and len(data_rows) == 2:
                row_min, row_max = int(data_rows[0]), int(data_rows[1])
            else:
                row_min, row_max = 2, sheet_data_rows + 1

            # Catégories (axe X / libellés)
            cat_col = chart_def.get("categories_col", 1)
            cats = Reference(
                src_ws,
                min_col=cat_col, max_col=cat_col,
                min_row=row_min, max_row=row_max,
            )

            # Séries
            series_defs = chart_def.get("series", [])
            if not series_defs:
                warnings.append(f"Graphique {idx+1} ('{chart_def.get('title', '')}') : aucune série définie, ignoré.")
                continue

            for s_def in series_defs:
                col = int(s_def.get("col", 2))
                vals = Reference(
                    src_ws,
                    min_col=col, max_col=col,
                    min_row=row_min, max_row=row_max,
                )

                if chart_type in ("scatter", "bubble"):
                    # ScatterChart / BubbleChart : x_values + values
                    x_vals = Reference(src_ws, min_col=cat_col, max_col=cat_col,
                                       min_row=row_min, max_row=row_max)
                    series = Series(vals, xvalues=x_vals)
                else:
                    series = Series(vals, cats)

                # Titre de la série (légende)
                if s_def.get("title"):
                    series.title = SeriesLabel(v=s_def["title"])

                # Couleur de remplissage
                if s_def.get("color"):
                    from openpyxl.drawing.fill import PatternFillProperties
                    from openpyxl.drawing.spreadsheet_drawing import SpreadsheetDrawing
                    try:
                        hex_color = s_def["color"].lstrip("#")
                        series.graphicalProperties.solidFill = hex_color
                    except Exception:
                        pass

                chart.series.append(series)

            # Étiquettes de données
            if chart_def.get("show_data_labels", False):
                try:
                    chart.dLbls = DataLabel()
                    chart.dLbls.showVal = True
                    chart.dLbls.showLegendKey = False
                    chart.dLbls.showCatName = False
                    chart.dLbls.showSerName = False
                except Exception:
                    pass

            # Cellule d'ancrage
            anchor = chart_def.get("anchor")
            if not anchor:
                # Par défaut : première colonne après les données, ligne 1
                next_col = (max((int(s.get("col", 2)) for s in series_defs), default=2) + 2)
                anchor = f"{get_column_letter(next_col)}1"

            ws.add_chart(chart, anchor)

        except Exception as e:
            warnings.append(f"Graphique {idx+1} ('{chart_def.get('title', '')}') : erreur — {e}")

    return warnings


@tool(
    name="export_xlsx_json",
    description=(
        "Génère un classeur Excel (.xlsx) depuis une structure JSON. "
        "Permet de créer plusieurs feuilles avec en-têtes, données et graphiques. "
        "Chaque feuille peut contenir un champ optionnel 'charts' pour générer "
        "des graphiques directement intégrés dans le classeur. "
        "Types de graphiques : bar, bar_stacked, bar_percent, bar_h, line, line_smooth, "
        "area, area_stacked, pie, doughnut, scatter, bubble, radar."
    ),
    parameters={
        "type": "object",
        "properties": {
            "workbook": {
                "type": "object",
                "description": (
                    "Structure du classeur. Champ 'sheets' : liste d'objets avec "
                    "name (str), headers (liste de str), rows (liste de listes), "
                    "et charts (liste optionnelle de graphiques). "
                    "Chaque graphique : {type, title, categories_col, series: [{title, col, color}], "
                    "data_rows, anchor, width_cm, height_cm, style, show_legend, show_data_labels}. "
                    "Exemple : voir la skill guide_export_excel.md pour la structure complète."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/données.xlsx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["workbook"]
    }
)
def export_xlsx_json(workbook: dict, output_path: str = "", filename: str = "") -> str:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter

        sheets = workbook.get("sheets", [])
        if not sheets:
            return _err("export_xlsx_json : 'sheets' est vide ou manquant")

        name = filename or "export.xlsx"
        if not name.endswith(".xlsx"):
            name += ".xlsx"
        p = _resolve_output(output_path, name)

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # supprimer la feuille vide par défaut

        header_font    = Font(bold=True, color="FFFFFF")
        header_fill    = PatternFill("solid", fgColor="2255A4")
        header_align   = Alignment(horizontal="center", vertical="center")
        alt_fill       = PatternFill("solid", fgColor="EFF3FA")

        total_rows  = 0
        all_warnings = []

        for sheet_def in sheets:
            ws = wb.create_sheet(title=str(sheet_def.get("name", "Feuille"))[:31])
            headers = sheet_def.get("headers", [])
            rows    = sheet_def.get("rows", [])

            # En-têtes
            for ci, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=ci, value=str(h))
                cell.font      = header_font
                cell.fill      = header_fill
                cell.alignment = header_align

            # Données + formatage alterné
            for ri, row in enumerate(rows, 2):
                fill = alt_fill if ri % 2 == 0 else None
                for ci, val in enumerate(row[:len(headers) or len(row)], 1):
                    cell = ws.cell(row=ri, column=ci, value=val)
                    if fill:
                        cell.fill = fill

            # Ajustement auto de la largeur des colonnes
            for col_cells in ws.columns:
                max_len = max((len(str(c.value or "")) for c in col_cells), default=8)
                ws.column_dimensions[get_column_letter(col_cells[0].column)].width = min(max_len + 4, 60)

            # Figer la première ligne
            ws.freeze_panes = "A2"
            total_rows += len(rows)

            # Graphiques
            charts_def = sheet_def.get("charts", [])
            if charts_def:
                warnings = _build_charts(wb, ws, charts_def, len(rows), len(headers))
                all_warnings.extend(warnings)

        wb.save(str(p))
        result = {"sheets": len(sheets), "total_rows": total_rows}
        if all_warnings:
            result["warnings"] = all_warnings
        return _ok(p, result)
    except Exception as e:
        return _err(f"export_xlsx_json : {e}")


@tool(
    name="export_xlsx_csv",
    description=(
        "Génère un fichier Excel (.xlsx) depuis du contenu CSV brut. "
        "Pratique quand les données sont déjà disponibles en format CSV "
        "(séparateur virgule ou point-virgule détecté automatiquement). "
        "Crée une seule feuille. Pour plusieurs feuilles, utiliser export_xlsx_json."
    ),
    parameters={
        "type": "object",
        "properties": {
            "csv_content": {
                "type": "string",
                "description": "Contenu CSV brut (avec en-têtes sur la première ligne)."
            },
            "sheet_name": {
                "type": "string",
                "description": "Nom de la feuille (défaut : 'Données')."
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination. Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["csv_content"]
    }
)
def export_xlsx_csv(csv_content: str, sheet_name: str = "Données",
                    output_path: str = "", filename: str = "") -> str:
    try:
        import csv as csv_mod
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter

        # Détecter le séparateur
        sample = csv_content[:2000]
        dialect = csv_mod.Sniffer().sniff(sample, delimiters=",;\t|")
        reader  = csv_mod.reader(io.StringIO(csv_content), dialect)
        all_rows = list(reader)

        if not all_rows:
            return _err("export_xlsx_csv : CSV vide")

        name = filename or "export.xlsx"
        if not name.endswith(".xlsx"):
            name += ".xlsx"
        p = _resolve_output(output_path, name)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = str(sheet_name)[:31]

        header_font  = Font(bold=True, color="FFFFFF")
        header_fill  = PatternFill("solid", fgColor="2255A4")
        header_align = Alignment(horizontal="center")
        alt_fill     = PatternFill("solid", fgColor="EFF3FA")

        for ri, row in enumerate(all_rows, 1):
            for ci, val in enumerate(row, 1):
                # Tenter une conversion numérique
                cell = ws.cell(row=ri, column=ci)
                try:
                    cell.value = int(val)
                except ValueError:
                    try:
                        cell.value = float(val.replace(",", "."))
                    except ValueError:
                        cell.value = val

                if ri == 1:
                    cell.font      = header_font
                    cell.fill      = header_fill
                    cell.alignment = header_align
                elif ri % 2 == 0:
                    cell.fill = alt_fill

        # Ajustement largeur
        for col_cells in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col_cells), default=8)
            ws.column_dimensions[get_column_letter(col_cells[0].column)].width = min(max_len + 4, 60)

        ws.freeze_panes = "A2"
        wb.save(str(p))
        return _ok(p, {"rows": len(all_rows) - 1, "columns": len(all_rows[0]) if all_rows else 0})
    except Exception as e:
        return _err(f"export_xlsx_csv : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# POWERPOINT / PPTX
# ═════════════════════════════════════════════════════════════════════════════

def _add_slide(prs, layout_idx: int, title: str, bullets: list[str],
               content: str = "", notes: str = "", subtitle: str = ""):
    """Ajoute un slide à la présentation."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    layout = prs.slide_layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)

    # Titre
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Sous-titre (layout 0 = titre principal)
    if layout_idx == 0 and subtitle:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = subtitle
                break

    # Corps / puces
    if bullets or content:
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx in (1, 2):  # corps du slide
                tf = ph.text_frame
                tf.clear()
                if bullets:
                    for i, b in enumerate(bullets):
                        if i == 0:
                            tf.paragraphs[0].text = b
                            tf.paragraphs[0].level = 0
                        else:
                            p = tf.add_paragraph()
                            p.text  = b
                            p.level = 0
                elif content:
                    tf.paragraphs[0].text = content
                break

    # Notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


@tool(
    name="export_pptx_json",
    description=(
        "Génère une présentation PowerPoint (.pptx) depuis une structure JSON. "
        "Chaque slide peut avoir un titre, une liste de puces, du contenu texte "
        "et des notes de présentation. "
        "À utiliser pour créer des présentations, des pitch decks, des diaporamas "
        "QUAND AUCUN GABARIT organisationnel n'est disponible. "
        "Si un fichier .pptx de modèle/gabarit existe dans ~/Modèles/ ou si "
        "l'utilisateur mentionne une charte graphique ou un modèle maison, "
        "utiliser export_pptx_template à la place. "
        "Avant tout export PowerPoint, appeler skill_list puis "
        "skill_read('guide_creation_powerpoint') pour respecter les conventions."
    ),
    parameters={
        "type": "object",
        "properties": {
            "presentation": {
                "type": "object",
                "description": (
                    "Structure de la présentation. Champs : "
                    "title (str, titre global), "
                    "subtitle (str, optionnel, pour le slide de titre), "
                    "slides (liste d'objets avec title/bullets/content/notes). "
                    "Exemple : {\"title\": \"Mon rapport\", \"slides\": ["
                    "{\"title\": \"Introduction\", \"bullets\": [\"Point 1\", \"Point 2\"], "
                    "\"notes\": \"Penser à mentionner le contexte.\"}]}"
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/présentation.pptx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["presentation"]
    }
)
def export_pptx_json(presentation: dict, output_path: str = "", filename: str = "") -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        name = filename or (presentation.get("title", "présentation") + ".pptx")
        if not name.endswith(".pptx"):
            name += ".pptx"
        p = _resolve_output(output_path, name)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide de titre
        main_title = presentation.get("title", "")
        subtitle   = presentation.get("subtitle", "")
        if main_title:
            _add_slide(prs, 0, main_title, [], subtitle=subtitle)

        # Slides de contenu
        for slide_def in presentation.get("slides", []):
            title   = slide_def.get("title", "")
            bullets = slide_def.get("bullets", [])
            content = slide_def.get("content", "")
            notes   = slide_def.get("notes", "")
            _add_slide(prs, 1, title, bullets, content=content, notes=notes)

        prs.save(str(p))
        n_slides = len(prs.slides)
        return _ok(p, {"slides": n_slides})
    except Exception as e:
        return _err(f"export_pptx_json : {e}")


@tool(
    name="export_pptx_outline",
    description=(
        "Génère une présentation PowerPoint (.pptx) depuis un outline texte simple. "
        "Format : '# Titre du slide' pour un nouveau slide, '- Puce' pour une puce, "
        "'> Note' pour une note de présentation, texte libre pour contenu. "
        "Plus rapide à écrire que le format JSON pour des présentations simples. "
        "À utiliser QUAND AUCUN GABARIT organisationnel n'est disponible. "
        "Si un fichier .pptx de modèle/gabarit existe dans ~/Modèles/ ou si "
        "l'utilisateur mentionne une charte graphique ou un modèle maison, "
        "utiliser export_pptx_template à la place. "
        "Avant tout export PowerPoint, appeler skill_list puis "
        "skill_read('guide_creation_powerpoint') pour respecter les conventions."
    ),
    parameters={
        "type": "object",
        "properties": {
            "outline": {
                "type": "string",
                "description": (
                    "Outline texte de la présentation. "
                    "Exemple :\\n"
                    "# Introduction\\n"
                    "- Contexte du projet\\n"
                    "- Objectifs\\n"
                    "> Penser à présenter l'équipe\\n"
                    "# Résultats\\n"
                    "- Hausse de 15% du CA\\n"
                )
            },
            "title": {
                "type": "string",
                "description": "Titre global de la présentation (slide de couverture)."
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination. Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["outline"]
    }
)
def export_pptx_outline(outline: str, title: str = "",
                         output_path: str = "", filename: str = "") -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches

        name = filename or (title or "présentation") + ".pptx"
        if not name.endswith(".pptx"):
            name += ".pptx"
        p = _resolve_output(output_path, name)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide de couverture
        if title:
            _add_slide(prs, 0, title, [])

        # Parser l'outline
        slides: list[dict] = []
        current: dict | None = None

        for raw_line in outline.splitlines():
            line = raw_line.rstrip()
            if line.startswith("# "):
                if current is not None:
                    slides.append(current)
                current = {"title": line[2:].strip(), "bullets": [], "notes": "", "content": ""}
            elif line.startswith("- ") and current is not None:
                current["bullets"].append(line[2:].strip())
            elif line.startswith("> ") and current is not None:
                current["notes"] += line[2:].strip() + " "
            elif line.strip() and current is not None and not current["bullets"]:
                current["content"] += line.strip() + " "

        if current is not None:
            slides.append(current)

        for s in slides:
            _add_slide(prs, 1, s["title"], s["bullets"],
                       content=s["content"].strip(), notes=s["notes"].strip())

        prs.save(str(p))
        return _ok(p, {"slides": len(prs.slides)})
    except Exception as e:
        return _err(f"export_pptx_outline : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# PDF
# ═════════════════════════════════════════════════════════════════════════════

# ── Rendu LaTeX → image PNG base64 (via latex + dvipng) ──────────────────────
#
# Pipeline :  expression LaTeX
#               → document .tex minimal (amsmath, amssymb, bm, preview, xcolor)
#               → latex → .dvi
#               → dvipng → .png transparent
#               → base64 data-URI  (PDF/HTML)  ou  bytes bruts  (DOCX)
#
# Avantages vs matplotlib.mathtext :
#   - Rendu TeX complet : \oint, \boldsymbol, \displaystyle, \bm, \text, etc.
#   - Pas de caractères ■ pour les symboles non supportés
#   - Qualité identique à un document LaTeX réel
#
# Prérequis système : latex (TeX Live) + dvipng
#   apt install texlive-latex-base texlive-latex-extra dvipng
#
# Note : _latex_to_png_bytes (section WORD) est la primitive commune.
#        _latex_to_base64_png (ci-dessous) la réutilise pour le pipeline PDF/HTML.

def _latex_to_base64_png(latex_expr: str, dpi: int = 200,
                          display: bool = False,
                          color: str = "1a1a2e") -> str:
    """
    Rend une expression LaTeX en data-URI PNG base64 pour injection HTML/CSS.

    Paramètres
    ----------
    latex_expr : expression LaTeX brute (sans les $ délimiteurs externes)
    dpi        : résolution (200 inline, 220 display)
    display    : ignoré — on utilise toujours equation* avec preview
    color      : couleur hex RGB sans # (défaut quasi-noir #1a1a2e)

    Délègue à _latex_to_png_bytes pour la compilation, puis encode en base64.
    Lève RuntimeError si latex ou dvipng échoue.
    """
    import base64

    # Si une couleur non standard est demandée, on passe par le chemin complet ;
    # sinon on délègue directement à _latex_to_png_bytes (couleur par défaut).
    if color.lower() in ("1a1a2e", "1a1a2e"):
        png_bytes = _latex_to_png_bytes(latex_expr, dpi=dpi)
    else:
        # Reconstruire avec couleur personnalisée — même double pipeline que _latex_to_png_bytes
        import shutil as _sh, subprocess as _sp, tempfile as _tmp
        from pathlib import Path as _P
        r_c = int(color[0:2], 16) / 255
        g_c = int(color[2:4], 16) / 255
        b_c = int(color[4:6], 16) / 255
        body = rf"\begin{{equation*}}{latex_expr}\end{{equation*}}"
        tex_src = rf"""\documentclass[12pt]{{article}}
\usepackage{{amsmath,amssymb,amsfonts,bm}}
\usepackage{{xcolor}}
\usepackage[active,tightpage]{{preview}}
\PreviewEnvironment{{equation*}}
\setlength\PreviewBorder{{3pt}}
\pagestyle{{empty}}
\begin{{document}}
\color[rgb]{{{r_c:.4f},{g_c:.4f},{b_c:.4f}}}
{body}
\end{{document}}
"""
        png_bytes = None
        with _tmp.TemporaryDirectory() as tmp:
            d = _P(tmp)
            (d / "f.tex").write_text(tex_src, encoding="utf-8")

            # Pipeline 1 : pdflatex + pdftoppm
            if _sh.which("pdflatex") and _sh.which("pdftoppm"):
                r1 = _sp.run(["pdflatex", "-interaction=nonstopmode", "-halt-on-error", "f.tex"],
                             cwd=d, capture_output=True, timeout=30)
                if (d / "f.pdf").exists():
                    _sp.run(["pdftoppm", "-r", str(dpi), "-png", "f.pdf", "page"],
                            cwd=d, capture_output=True, timeout=20)
                    pages = sorted(d.glob("page*.png"))
                    if pages:
                        raw = pages[0].read_bytes()
                        try:
                            import io as _io2, numpy as _np
                            from PIL import Image as _Img
                            img = _Img.open(_io2.BytesIO(raw)).convert("RGBA")
                            arr = _np.array(img)
                            white = (arr[:,:,0] > 250) & (arr[:,:,1] > 250) & (arr[:,:,2] > 250)
                            arr[white, 3] = 0
                            buf = _io2.BytesIO()
                            _Img.fromarray(arr, "RGBA").save(buf, format="PNG")
                            png_bytes = buf.getvalue()
                        except Exception:
                            png_bytes = raw

            # Pipeline 2 : latex + dvipng (repli)
            if png_bytes is None and _sh.which("dvipng"):
                r1 = _sp.run(["latex", "-interaction=nonstopmode", "-halt-on-error", "f.tex"],
                             cwd=d, capture_output=True, timeout=20)
                if (d / "f.dvi").exists():
                    _sp.run(["dvipng", "-D", str(dpi), "-T", "tight", "--png",
                             "-bg", "Transparent", "-o", "f.png", "f.dvi"],
                            cwd=d, capture_output=True, timeout=20)
                    if (d / "f.png").exists():
                        png_bytes = (d / "f.png").read_bytes()

        if png_bytes is None:
            raise RuntimeError(f"_latex_to_base64_png : aucun pipeline disponible")

    return "data:image/png;base64," + base64.b64encode(png_bytes).decode()


def _render_latex_in_text(text: str) -> str:
    """
    Convertit un texte avec Markdown inline, blocs de code et formules LaTeX en HTML.

    Marqueurs Markdown gérés :
      **gras**         → <strong>
      *italique*       → <em>
      __souligné__     → <u>
      `code inline`    → <code> monospace inline avec fond gris
      ```lang\\n…\\n``` → <pre><code> bloc multilignes avec bandeau langue coloré

    Retours à la ligne dans les segments texte :
      \\n → <br> (saut de ligne visible dans le PDF)

    Formules LaTeX :
      $...$   → <img> inline, height:1.4em
      $$...$$ → <div> centré, hauteur calculée depuis les dims réelles du PNG

    En cas d'échec du rendu LaTeX, un <code> de repli rouge est inséré.
    """
    import html as _html
    import base64 as _b64

    def _png_src(png_bytes: bytes) -> str:
        return "data:image/png;base64," + _b64.b64encode(png_bytes).decode()

    segs = _parse_inline_markup(text)
    parts: list[str] = []

    for s in segs:
        t = s["type"]

        if t == "text":
            # \n → <br> pour préserver les sauts de ligne intentionnels
            escaped = _html.escape(s["content"])
            parts.append(escaped.replace("\n", "<br>"))

        elif t == "bold":
            parts.append(f'<strong>{_html.escape(s["content"])}</strong>')

        elif t == "italic":
            parts.append(f'<em>{_html.escape(s["content"])}</em>')

        elif t == "underline":
            parts.append(f'<u>{_html.escape(s["content"])}</u>')

        elif t == "code":
            parts.append(
                f'<code style="font-family:\'Courier New\',Courier,monospace;'
                f'font-size:0.88em;background:#F4F6FA;'
                f'padding:1px 4px;border-radius:3px;'
                f'border:0.5pt solid #C5CDE0;">'
                f'{_html.escape(s["content"])}</code>'
            )

        elif t == "code_block":
            lang    = s["lang"]
            content = _html.escape(s["content"].rstrip("\n"))
            # Bandeau coloré avec libellé du langage (si précisé)
            lang_badge = (
                f'<span style="font-family:\'Courier New\',Courier,monospace;'
                f'font-size:7.5pt;color:#ffffff;float:right;padding:2pt 8pt;">'
                f'{_html.escape(lang)}</span>'
            ) if lang else ""
            parts.append(
                f'<div style="margin:8pt 0;border:0.5pt solid #C5CDE0;'
                f'border-radius:4pt;overflow:hidden;page-break-inside:avoid;">'
                f'<div style="background:#2255A4;padding:3pt 8pt;'
                f'min-height:14pt;line-height:14pt;">{lang_badge}</div>'
                f'<pre style="margin:0;padding:8pt 10pt;background:#F8F9FB;'
                f'font-family:\'Courier New\',Courier,monospace;font-size:8.5pt;'
                f'line-height:1.55;white-space:pre;color:#1a1a2e;">'
                f'<code>{content}</code></pre></div>'
            )

        elif t == "latex_inline":
            try:
                png = _latex_to_png_bytes(s["expr"], dpi=200)
                src = _png_src(png)
                alt = _html.escape(s["expr"])
                parts.append(
                    f'<img src="{src}" alt="{alt}" '
                    f'style="height:1.4em;vertical-align:middle;'
                    f'display:inline-block;margin:0 1px;">'
                )
            except Exception as exc:
                parts.append(
                    f'<code style="color:#c0392b;" title="{_html.escape(str(exc))}">'
                    f'{_html.escape(s["expr"])}</code>'
                )

        elif t == "latex_display":
            try:
                png = _latex_to_png_bytes(s["expr"], dpi=220)
                src = _png_src(png)
                alt = _html.escape(s["expr"])
                _, h_cm = _latex_png_dims(png, 220)
                h_pt = min(h_cm * 28.35, 60.0)
                parts.append(
                    f'<div style="text-align:center;margin:8pt 0 6pt;">'
                    f'<img src="{src}" alt="{alt}" '
                    f'style="height:{h_pt:.1f}pt;max-width:90%;vertical-align:middle;">'
                    f'</div>'
                )
            except Exception as exc:
                parts.append(
                    f'<code style="color:#c0392b;" title="{_html.escape(str(exc))}">'
                    f'{_html.escape(s["expr"])}</code>'
                )

    return "".join(parts)


def _doc_to_html(document: dict) -> str:
    """
    Convertit la structure JSON commune (title / sections) en HTML complet,
    avec :
      - rendu LaTeX inline/display via matplotlib
      - mise en forme CSS professionnelle (palette bleue Prométhée)
      - en-tête, pied de page, numérotation automatique WeasyPrint (@page)
      - tableaux, listes à puces, sauts de page
    """
    import html as html_mod

    BLUE   = "#2255A4"
    BLUE2  = "#3B7ACC"
    GREY   = "#F4F6FA"
    DGREY  = "#374151"
    LGREY  = "#C5CDE0"

    title  = document.get("title", "")
    footer_label = html_mod.escape(title or "Document généré par Prométhée AI")

    CSS = f"""
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700&display=swap');

    @page {{
        size: A4;
        margin: 2cm 1.5cm 2.5cm 1.5cm;

        @bottom-left {{
            content: "{footer_label}";
            font-family: 'Inter', Helvetica, Arial, sans-serif;
            font-size: 7.5pt;
            color: #6B7280;
        }}
        @bottom-right {{
            content: "Page " counter(page);
            font-family: 'Inter', Helvetica, Arial, sans-serif;
            font-size: 7.5pt;
            color: #6B7280;
        }}
    }}

    * {{ box-sizing: border-box; margin: 0; padding: 0; }}

    body {{
        font-family: 'Inter', Helvetica, Arial, sans-serif;
        font-size: 10pt;
        color: {DGREY};
        line-height: 1.6;
        background: white;
    }}

    h1.doc-title {{
        font-size: 18pt;
        color: {BLUE};
        font-weight: 700;
        margin-bottom: 6pt;
        hyphens: none;
    }}
    hr.title-rule {{
        border: none;
        border-top: 2px solid {BLUE};
        margin-bottom: 14pt;
    }}

    h2 {{ font-size: 14pt; color: {BLUE};  font-weight: 700; margin: 18pt 0 6pt; hyphens: none; }}
    h3 {{ font-size: 12pt; color: {BLUE2}; font-weight: 600; margin: 12pt 0 4pt; hyphens: none; }}
    h4 {{ font-size: 11pt; color: {DGREY}; font-weight: 600; margin:  8pt 0 3pt; hyphens: none; }}

    p {{
        text-align: justify;
        hyphens: auto;
        hyphenate-limit-chars: 6 3 3;
        hyphenate-limit-lines: 2;
        word-spacing: -0.01em;
        margin-bottom: 6pt;
        orphans: 3;
        widows: 3;
    }}

    ul.blt {{
        list-style: none;
        padding-left: 14pt;
        margin-bottom: 6pt;
    }}
    ul.blt li::before {{
        content: "▸ ";
        color: {BLUE};
        font-weight: 700;
    }}
    ul.blt li {{
        margin-bottom: 3pt;
        line-height: 1.5;
        text-align: left;
        hyphens: none;
    }}

    table {{
        width: 100%;
        border-collapse: collapse;
        margin-bottom: 10pt;
        font-size: 9pt;
    }}
    thead tr {{
        background: {BLUE};
    }}
    thead th {{
        color: white;
        font-weight: 700;
        padding: 5pt 6pt;
        text-align: left;
        border: 0.4pt solid {LGREY};
        hyphens: none;
    }}
    tbody tr:nth-child(odd)  {{ background: white; }}
    tbody tr:nth-child(even) {{ background: {GREY}; }}
    tbody td {{
        padding: 5pt 6pt;
        border: 0.4pt solid {LGREY};
        vertical-align: top;
        hyphens: auto;
    }}

    .page-break {{ page-break-after: always; }}

    img.latex-display {{
        display: block;
        margin: 10pt auto;
        max-width: 95%;
    }}
    """

    def _esc_latex(text: str) -> str:
        """
        Convertit un texte avec Markdown inline et formules LaTeX en HTML.
        Délègue entièrement à _render_latex_in_text qui gère l'échappement HTML,
        le rendu LaTeX (PNG base64) et les marqueurs **gras**, *italique*, etc.
        """
        return _render_latex_in_text(str(text))

    body_parts: list[str] = []

    if title:
        body_parts.append(f'<h1 class="doc-title">{html_mod.escape(title)}</h1>')
        body_parts.append('<hr class="title-rule">')

    heading_tags = {1: "h2", 2: "h3", 3: "h4"}

    for sec in document.get("sections", []):
        heading = sec.get("heading")
        level   = max(1, min(3, int(sec.get("level", 1))))

        if heading:
            tag = heading_tags[level]
            body_parts.append(f"<{tag}>{_esc_latex(heading)}</{tag}>")

        # ── Paragraphes ───────────────────────────────────────────────────
        if sec.get("paragraphs"):
            for para in sec["paragraphs"]:
                if para and str(para).strip():
                    body_parts.append(f"<p>{_esc_latex(str(para))}</p>")
        elif sec.get("content"):
            raw = sec["content"]
            if "\n\n" in raw:
                parts_txt = [x.strip() for x in raw.split("\n\n") if x.strip()]
            else:
                parts_txt = [x.strip() for x in raw.split("\n") if x.strip()]
            for part in parts_txt:
                body_parts.append(f"<p>{_esc_latex(part)}</p>")

        # ── Intro ────────────────────────────────────────────────────────
        if sec.get("intro"):
            body_parts.append(f"<p>{_esc_latex(str(sec['intro']))}</p>")

        # ── Liste à puces ─────────────────────────────────────────────────
        if sec.get("bullets"):
            items = "".join(
                f"<li>{_esc_latex(str(b))}</li>"
                for b in sec["bullets"] if b and str(b).strip()
            )
            body_parts.append(f'<ul class="blt">{items}</ul>')

        # ── Tableau ───────────────────────────────────────────────────────
        if sec.get("table"):
            tbl_data = sec["table"]
            headers  = tbl_data.get("headers", [])
            rows     = tbl_data.get("rows", [])
            if headers:
                ths = "".join(f"<th>{_esc_latex(str(h))}</th>" for h in headers)
                thead = f"<thead><tr>{ths}</tr></thead>"
                trs = ""
                for row in rows:
                    tds = "".join(
                        f"<td>{_esc_latex(str(v))}</td>"
                        for v in row[:len(headers)]
                    )
                    trs += f"<tr>{tds}</tr>"
                body_parts.append(f"<table>{thead}<tbody>{trs}</tbody></table>")

        # ── Saut de page ──────────────────────────────────────────────────
        if sec.get("page_break"):
            body_parts.append('<div class="page-break"></div>')

    body_html = "\n".join(body_parts)

    return f"""<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="utf-8">
<title>{html_mod.escape(title)}</title>
<style>{CSS}</style>
</head>
<body>
{body_html}
</body>
</html>"""


@tool(
    name="export_pdf",
    description=(
        "Génère un document PDF structuré depuis une description JSON. "
        "Supporte les titres hiérarchiques (niveaux 1 à 3), les paragraphes, "
        "les listes à puces, les tableaux et les formules mathématiques LaTeX. "
        "Les formules LaTeX inline ($...$) et display ($$...$$) sont rendu "
        "automatiquement en images vectorielles via matplotlib. "
        "Le moteur de rendu est WeasyPrint (HTML/CSS → PDF) : mise en page "
        "professionnelle avec numérotation des pages et pied de page automatiques. "
        "À utiliser pour produire des rapports, fiches, documents formels, "
        "documents scientifiques ou techniques contenant des formules mathématiques. "
        "Même format JSON que export_docx : utiliser 'paragraphs' (liste) pour plusieurs "
        "paragraphes, 'intro' avant les bullets/tableau, et rédiger un contenu dense "
        "(15 à 40 sections pour un document professionnel complet)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "document": {
                "type": "object",
                "description": (
                    "Structure du document (même format que export_docx). "
                    "Champs : title (str), sections (liste d'objets avec "
                    "heading, level, paragraphs, content, intro, bullets, table, page_break). "
                    "Les formules LaTeX peuvent être intégrées dans tout champ textuel : "
                    "inline avec $formule$ (ex: $E=mc^2$) "
                    "ou en bloc centré avec $$formule$$ (ex: $$\\int_0^\\infty e^{-x}dx=1$$). "
                    "Rédiger un contenu complet et développé."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/rapport.pdf). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["document"]
    }
)
def export_pdf(document: dict, output_path: str = "", filename: str = "") -> str:
    """
    Export PDF via WeasyPrint (HTML/CSS → PDF).
    Pipeline :
      1. _doc_to_html()    : JSON → HTML (avec formules LaTeX converties en PNG base64)
      2. weasyprint        : HTML → PDF avec @page, pied de page, numérotation
    Repli automatique sur reportlab si WeasyPrint n'est pas disponible.
    """
    name = filename or (document.get("title", "export") + ".pdf")
    if not name.endswith(".pdf"):
        name += ".pdf"
    p = _resolve_output(output_path, name)

    # ── Tentative WeasyPrint ───────────────────────────────────────────────────
    try:
        import weasyprint  # noqa: F401 — test de disponibilité
        html_source = _doc_to_html(document)
        weasyprint.HTML(string=html_source).write_pdf(str(p))
        return _ok(p, {"sections": len(document.get("sections", [])),
                        "engine": "weasyprint"})

    except ImportError:
        pass   # WeasyPrint absent → repli reportlab ci-dessous
    except Exception as e:
        return _err(f"export_pdf (weasyprint) : {e}")

    # ── Repli reportlab (sans LaTeX) ──────────────────────────────────────────
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib.enums import TA_JUSTIFY
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, HRFlowable,
        )

        BLUE  = colors.HexColor("#2255A4")
        GREY  = colors.HexColor("#F4F6FA")
        GLINE = colors.HexColor("#C5CDE0")

        def ST(n, **kw): return ParagraphStyle(n, **kw)
        styles = {
            "h0":   ST("h0",   fontName="Helvetica-Bold", fontSize=18, textColor=BLUE,
                        spaceBefore=0, spaceAfter=10),
            "h1":   ST("h1",   fontName="Helvetica-Bold", fontSize=14, textColor=BLUE,
                        spaceBefore=18, spaceAfter=6),
            "h2":   ST("h2",   fontName="Helvetica-Bold", fontSize=12,
                        textColor=colors.HexColor("#3B7ACC"), spaceBefore=12, spaceAfter=4),
            "h3":   ST("h3",   fontName="Helvetica-Bold", fontSize=11,
                        textColor=colors.HexColor("#374151"), spaceBefore=8, spaceAfter=3),
            "body": ST("body", fontName="Helvetica", fontSize=10,
                        leading=15, spaceBefore=2, spaceAfter=6, alignment=TA_JUSTIFY),
            "blt":  ST("blt",  fontName="Helvetica", fontSize=10, leading=14,
                        spaceBefore=2, spaceAfter=2, leftIndent=14),
            "th":   ST("th",   fontName="Helvetica-Bold", fontSize=9,
                        textColor=colors.white, leading=11),
            "tc":   ST("tc",   fontName="Helvetica", fontSize=9,
                        textColor=colors.HexColor("#111827"), leading=12),
        }

        def footer(canvas, doc):
            canvas.saveState()
            canvas.setFont("Helvetica", 7.5)
            canvas.setFillColor(colors.HexColor("#6B7280"))
            canvas.drawString(1.5*cm, 1*cm,
                              document.get("title", "Document généré par Prométhée AI"))
            canvas.drawRightString(A4[0] - 1.5*cm, 1*cm, f"Page {doc.page}")
            canvas.restoreState()

        story = []
        W = A4[0] - 3*cm

        if document.get("title"):
            story.append(Paragraph(document["title"], styles["h0"]))
            story.append(HRFlowable(width="100%", thickness=2, color=BLUE, spaceAfter=8))

        heading_map = {1: "h1", 2: "h2", 3: "h3"}

        def _add_para(txt: str):
            """
            Ajoute un paragraphe ReportLab en gérant les formules LaTeX.
            - Texte pur → Paragraph simple.
            - Formule display $$...$$ → Image centrée (PNG via latex+dvipng).
            - Formule inline $...$ → Image inline insérée via <img> ReportLab
              (platypus ne supporte pas les images inline dans Paragraph ;
               on éclate le texte autour de chaque formule et on insère des
               flowables séquentiels : Paragraph / Image centré / Paragraph).
            En cas d'échec du rendu, on insère le texte brut entre crochets.
            """
            import io as _io
            from reportlab.platypus import Image as _RLImage
            from reportlab.lib.units import pt as _pt

            segs = _parse_inline_markup(txt)
            has_display = any(s["type"] == "latex_display" for s in segs)
            has_inline  = any(s["type"] == "latex_inline"  for s in segs)

            if not has_display and not has_inline:
                # Pas de LaTeX : paragraphe simple
                story.append(Paragraph(txt, styles["body"]))
                return

            # Accumuler les fragments texte courants ; à chaque formule display
            # on flush le texte accumulé puis on insère l'image centrée.
            # Pour les formules inline on les encode comme [expr] en fallback
            # ou on génère une image séparée centrée (ReportLab ne supporte pas
            # les images inline dans Paragraph sans extension tiers).
            pending_html: list[str] = []

            def _flush_text():
                joined = "".join(pending_html).strip()
                if joined:
                    story.append(Paragraph(joined, styles["body"]))
                pending_html.clear()

            for seg in segs:
                t = seg["type"]
                if t == "text":
                    import html as _h
                    pending_html.append(_h.escape(seg["content"]).replace("\n", "<br/>"))
                elif t == "bold":
                    import html as _h
                    pending_html.append(f'<b>{_h.escape(seg["content"])}</b>')
                elif t == "italic":
                    import html as _h
                    pending_html.append(f'<i>{_h.escape(seg["content"])}</i>')
                elif t == "underline":
                    import html as _h
                    pending_html.append(f'<u>{_h.escape(seg["content"])}</u>')
                elif t == "code":
                    import html as _h
                    pending_html.append(
                        f'<font name="Courier" size="9">{_h.escape(seg["content"])}</font>'
                    )
                elif t == "latex_inline":
                    # Flush le texte avant, puis insérer la formule centrée
                    _flush_text()
                    try:
                        png = _latex_to_png_bytes(seg["expr"], dpi=200)
                        w_cm, h_cm = _latex_png_dims(png, 200)
                        h_pt = min(h_cm * 28.35, 14.0)
                        w_pt = w_cm * 28.35 * (h_pt / (h_cm * 28.35))
                        img = _RLImage(_io.BytesIO(png), width=w_pt, height=h_pt)
                        img.hAlign = "CENTER"
                        story.append(img)
                    except Exception:
                        story.append(Paragraph(
                            f'<font color="#c0392b">[{seg["expr"]}]</font>',
                            styles["body"]
                        ))
                elif t == "latex_display":
                    _flush_text()
                    try:
                        png = _latex_to_png_bytes(seg["expr"], dpi=220)
                        w_cm, h_cm = _latex_png_dims(png, 220)
                        max_w_pt = W
                        h_pt = min(h_cm * 28.35, 60.0)
                        w_pt = min(w_cm * 28.35, max_w_pt)
                        # Conserver le ratio
                        ratio = w_cm / h_cm if h_cm else 1
                        if w_pt / h_pt > ratio:
                            w_pt = h_pt * ratio
                        img = _RLImage(_io.BytesIO(png), width=w_pt, height=h_pt)
                        img.hAlign = "CENTER"
                        story.append(img)
                    except Exception:
                        story.append(Paragraph(
                            f'<font color="#c0392b">[{seg["expr"]}]</font>',
                            styles["body"]
                        ))

            _flush_text()

        for sec in document.get("sections", []):
            heading = sec.get("heading")
            level   = max(1, min(3, int(sec.get("level", 1))))
            if heading:
                story.append(Paragraph(heading, styles[heading_map[level]]))

            if sec.get("paragraphs"):
                for para in sec["paragraphs"]:
                    if para and str(para).strip():
                        _add_para(str(para))
            elif sec.get("content"):
                raw = sec["content"]
                parts_rl = ([x.strip() for x in raw.split("\n\n") if x.strip()]
                            if "\n\n" in raw
                            else [x.strip() for x in raw.split("\n") if x.strip()])
                for part in parts_rl:
                    _add_para(part)

            if sec.get("intro"):
                _add_para(str(sec["intro"]))

            if sec.get("bullets"):
                for b in sec["bullets"]:
                    if b and str(b).strip():
                        story.append(Paragraph(
                            f'<font color="#2255A4">&#8226;</font>  {b}',
                            styles["blt"]
                        ))
                story.append(Spacer(1, 4))

            if sec.get("table"):
                tbl_data = sec["table"]
                headers  = tbl_data.get("headers", [])
                rows     = tbl_data.get("rows", [])
                if headers:
                    ncols = len(headers)
                    col_w = W / ncols
                    data  = [[Paragraph(h, styles["th"]) for h in headers]]
                    for row in rows:
                        data.append([Paragraph(str(v), styles["tc"]) for v in row[:ncols]])
                    ts = TableStyle([
                        ("BACKGROUND",    (0, 0), (-1, 0),  BLUE),
                        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, GREY]),
                        ("GRID",          (0, 0), (-1, -1), 0.4, GLINE),
                        ("VALIGN",        (0, 0), (-1, -1), "TOP"),
                        ("TOPPADDING",    (0, 0), (-1, -1), 5),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING",  (0, 0), (-1, -1), 6),
                    ])
                    story.append(Table(data, colWidths=[col_w]*ncols, repeatRows=1, style=ts))
                    story.append(Spacer(1, 8))

            if sec.get("page_break"):
                story.append(PageBreak())

        doc_obj = SimpleDocTemplate(
            str(p), pagesize=A4,
            topMargin=2*cm, bottomMargin=2.2*cm,
            leftMargin=1.5*cm, rightMargin=1.5*cm,
            title=document.get("title", ""),
        )
        doc_obj.build(story, onFirstPage=footer, onLaterPages=footer)
        return _ok(p, {"sections": len(document.get("sections", [])),
                        "engine": "reportlab (repli — weasyprint absent)"})

    except Exception as e:
        return _err(f"export_pdf (reportlab fallback) : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# LIBREOFFICE — CONVERSION ET NATIF
# ═════════════════════════════════════════════════════════════════════════════

def _libreoffice_convert(input_path: Path, target_format: str,
                          output_dir: Path) -> Path | None:
    """
    Appelle LibreOffice headless pour convertir un fichier.
    Retourne le chemin du fichier converti ou None en cas d'erreur.
    """
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        return None

    cmd = [
        soffice, "--headless", "--norestore",
        f"--convert-to", target_format,
        "--outdir", str(output_dir),
        str(input_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice : {result.stderr.strip()}")

    # LibreOffice place le fichier converti dans output_dir
    # avec le même stem que l'entrée
    ext_map = {
        "odt": ".odt", "ods": ".ods", "odp": ".odp",
        "pdf": ".pdf", "docx": ".docx", "xlsx": ".xlsx",
    }
    stem = input_path.stem
    ext  = ext_map.get(target_format, f".{target_format}")
    converted = output_dir / (stem + ext)
    if converted.exists():
        return converted

    # Chercher tout fichier avec la bonne extension (LibreOffice peut varier)
    candidates = list(output_dir.glob(f"{stem}*{ext}"))
    return candidates[0] if candidates else None


@tool(
    name="export_libreoffice",
    description=(
        "Convertit un fichier existant (docx, xlsx, pptx, pdf, csv…) vers un format "
        "LibreOffice natif : .odt (traitement de texte), .ods (tableur), .odp (présentation). "
        "Utilise LibreOffice en mode headless. Le fichier source doit exister sur le disque."
    ),
    parameters={
        "type": "object",
        "properties": {
            "input_path": {
                "type": "string",
                "description": "Chemin absolu ou relatif au home du fichier à convertir."
            },
            "target_format": {
                "type": "string",
                "enum": ["odt", "ods", "odp"],
                "description": "Format LibreOffice cible : odt (texte), ods (tableur), odp (présentation)."
            },
            "output_path": {
                "type": "string",
                "description": "Chemin du fichier de sortie. Si omis, placé à côté du fichier source."
            }
        },
        "required": ["input_path", "target_format"]
    }
)
def export_libreoffice(input_path: str, target_format: str,
                        output_path: str = "") -> str:
    try:
        src = Path(input_path).expanduser().resolve()
        if not src.exists():
            return _err(f"export_libreoffice : fichier source introuvable : {src}")

        # Dossier de sortie temporaire pour LibreOffice
        with tempfile.TemporaryDirectory() as tmp_dir:
            converted = _libreoffice_convert(src, target_format, Path(tmp_dir))
            if converted is None:
                return _err("export_libreoffice : LibreOffice non disponible ou conversion échouée")

            # Déterminer la destination finale
            if output_path:
                dest = Path(output_path).expanduser()
                if not dest.is_absolute():
                    dest = Path.home() / dest
            else:
                dest = src.parent / (src.stem + f".{target_format}")

            dest.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(str(converted), str(dest))

        return _ok(dest)
    except Exception as e:
        return _err(f"export_libreoffice : {e}")


@tool(
    name="export_libreoffice_native",
    description=(
        "Génère directement un fichier LibreOffice natif (.odt, .ods ou .odp) "
        "depuis une description JSON, sans passer par Word/Excel/PowerPoint. "
        "Le document est d'abord construit dans le format intermédiaire le plus "
        "approprié (docx→odt, xlsx→ods, pptx→odp) puis converti via LibreOffice. "
        "Même format JSON que export_docx / export_xlsx_json / export_pptx_json."
    ),
    parameters={
        "type": "object",
        "properties": {
            "target_format": {
                "type": "string",
                "enum": ["odt", "ods", "odp"],
                "description": "Format de sortie LibreOffice : odt (texte), ods (tableur), odp (présentation)."
            },
            "document": {
                "type": "object",
                "description": (
                    "Structure du document. Pour odt : même format que export_docx. "
                    "Pour ods : même format que export_xlsx_json (champ 'sheets'). "
                    "Pour odp : même format que export_pptx_json (champ 'slides')."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination du fichier LibreOffice final. Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["target_format", "document"]
    }
)
def export_libreoffice_native(target_format: str, document: dict,
                               output_path: str = "", filename: str = "") -> str:
    try:
        doc_title = document.get("title", "export")
        name = filename or f"{doc_title}.{target_format}"
        if not name.endswith(f".{target_format}"):
            name = name.rsplit(".", 1)[0] + f".{target_format}"
        final_dest = _resolve_output(output_path, name)

        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp = Path(tmp_dir)

            # ── Étape 1 : générer le format intermédiaire ──────────────────
            if target_format == "odt":
                tmp_src = tmp / f"{doc_title}.docx"
                doc = _build_docx(document)
                doc.save(str(tmp_src))
                intermediate_format = "odt"

            elif target_format == "ods":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                from openpyxl.utils import get_column_letter

                tmp_src = tmp / f"{doc_title}.xlsx"
                wb = openpyxl.Workbook()
                wb.remove(wb.active)
                hf = Font(bold=True, color="FFFFFF")
                hfill = PatternFill("solid", fgColor="2255A4")
                ha = Alignment(horizontal="center")

                for sheet_def in document.get("sheets", []):
                    ws = wb.create_sheet(str(sheet_def.get("name", "Feuille"))[:31])
                    headers = sheet_def.get("headers", [])
                    rows    = sheet_def.get("rows", [])
                    for ci, h in enumerate(headers, 1):
                        c = ws.cell(row=1, column=ci, value=str(h))
                        c.font = hf; c.fill = hfill; c.alignment = ha
                    for ri, row in enumerate(rows, 2):
                        for ci, val in enumerate(row, 1):
                            ws.cell(row=ri, column=ci, value=val)
                    for col_cells in ws.columns:
                        max_len = max((len(str(c.value or "")) for c in col_cells), default=8)
                        ws.column_dimensions[get_column_letter(col_cells[0].column)].width = min(max_len + 4, 60)
                    ws.freeze_panes = "A2"
                wb.save(str(tmp_src))
                intermediate_format = "ods"

            elif target_format == "odp":
                from pptx import Presentation
                from pptx.util import Inches
                tmp_src = tmp / f"{doc_title}.pptx"
                prs = Presentation()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                main_title = document.get("title", "")
                subtitle   = document.get("subtitle", "")
                if main_title:
                    _add_slide(prs, 0, main_title, [], subtitle=subtitle)
                for s in document.get("slides", []):
                    _add_slide(prs, 1, s.get("title", ""), s.get("bullets", []),
                               content=s.get("content", ""), notes=s.get("notes", ""))
                prs.save(str(tmp_src))
                intermediate_format = "odp"
            else:
                return _err(f"export_libreoffice_native : format non supporté : {target_format}")

            # ── Étape 2 : conversion via LibreOffice ───────────────────────
            converted = _libreoffice_convert(tmp_src, intermediate_format, tmp)
            if converted is None:
                return _err(
                    "export_libreoffice_native : LibreOffice non disponible. "
                    f"Le fichier intermédiaire est disponible : {tmp_src}"
                )

            final_dest.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(str(converted), str(final_dest))

        return _ok(final_dest, {"intermediate_format": intermediate_format})
    except Exception as e:
        return _err(f"export_libreoffice_native : {e}")
